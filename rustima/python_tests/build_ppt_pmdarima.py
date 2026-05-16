#!/usr/bin/env python3
"""
rustima vs pmdarima 벤치마크 결과 PPT 생성 (3장).
실행: .venv/bin/python python_tests/build_ppt_pmdarima.py
산출물: paper/bench_2019_2023/rustima_vs_pmdarima.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
BENCH_DIR = os.path.join(ROOT, "..", "paper", "bench_2019_2023")
CHART_COMPARE = os.path.join(BENCH_DIR, "compare_pmdarima_1y.png")
CHART_SCALING = os.path.join(BENCH_DIR, "auto_arima_scaling.png")
CHART_UNIFIED = os.path.join(BENCH_DIR, "scaling_with_pmdarima.png")
OUT_PATH = os.path.join(BENCH_DIR, "rustima_vs_pmdarima.pptx")

# Palette
ORANGE = RGBColor(0xDE, 0x6B, 0x35)
DARK   = RGBColor(0x22, 0x22, 0x33)
GRAY   = RGBColor(0x66, 0x66, 0x77)
LIGHT  = RGBColor(0xEE, 0xEE, 0xF2)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x30, 0x6E, 0xA8)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

FONT = "Pretendard"  # falls back if missing


# ───── helpers ─────────────────────────────────────────────────────────────
def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, font=FONT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    return tb


def add_title(slide, title, subtitle=None):
    add_text(slide, 0.5, 0.25, 12.3, 0.6, title,
             size=28, bold=True, color=DARK)
    if subtitle:
        add_text(slide, 0.5, 0.85, 12.3, 0.45, subtitle,
                 size=14, color=GRAY)
    # Orange accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.32), Inches(0.7), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()


def add_bullet(slide, x, y, w, h, items, size=14, line_spacing=1.4):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, txt in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + txt
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.font.name = FONT
        p.line_spacing = line_spacing
    return tb


def add_image(slide, path, x, y, w, h):
    return slide.shapes.add_picture(path, Inches(x), Inches(y),
                                    width=Inches(w), height=Inches(h))


def add_table(slide, x, y, w, h, data, header=True,
              first_col_bold=False, highlight_row=None, highlight_col=None):
    rows, cols = len(data), len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y),
                                       Inches(w), Inches(h))
    tbl = tbl_shape.table
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c])
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in para.runs:
                    run.font.size = Pt(13)
                    run.font.name = FONT
                    if header and r == 0:
                        run.font.bold = True
                        run.font.color.rgb = WHITE
                    elif first_col_bold and c == 0:
                        run.font.bold = True
                        run.font.color.rgb = DARK
                    else:
                        run.font.color.rgb = DARK
            # backgrounds
            if header and r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK
            elif highlight_row is not None and r == highlight_row:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xEE)
            elif highlight_col is not None and c == highlight_col:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xEE)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return tbl_shape


def add_callout(slide, x, y, w, h, text, color=ORANGE, size=18):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xEE)
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.15)
    tf.margin_top = tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = FONT


def add_footer(slide, text):
    add_text(slide, 0.5, 7.05, 12.3, 0.35, text,
             size=10, color=GRAY, align=PP_ALIGN.RIGHT)


# ───── slide 1: unified scaling + comparison ──────────────────────────────
def slide_performance():
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_title(s, "rustima vs pmdarima — 동일 축 위 직접 비교",
              "rustima는 1y → 5y 전 구간 측정 · pmdarima는 1y만 측정 (다년치는 OOM 위험)")

    # full-width unified chart
    add_image(s, CHART_UNIFIED, x=0.5, y=1.55, w=12.3, h=4.4)

    # bottom: key numbers strip
    add_bullet(s, 0.7, 6.05, 12.0, 1.0, [
        "1년치 직접 비교 — wall time 76.2 s vs 404.0 s (5.3× 빠름) · peak RSS 248 MB vs 8,685 MB (35× 적음)",
        "rustima 확장성 — 5년치(43,824 obs)도 322 s · 638 MB로 완료 (시간 ≈ 7.3 ms/obs)",
        "pmdarima는 1년치만으로 8.7 GB → 다년치는 측정 불가 (그래프 빗금 영역)",
    ], size=13)

    add_footer(s, "Apple M-series · stepwise auto_arima · max_p=q=3, max_P=Q=1, max_d=D=1")


# ───── slide 2: model quality + conclusion ────────────────────────────────
def slide_quality():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "빠를 뿐 아니라 더 잘 맞는다",
              "선택된 모델의 적합도 비교 (auto_arima 결과)")

    # left: comparison table
    table_data = [
        ["지표", "rustima", "pmdarima"],
        ["차수", "(3,0,3)(1,1,1)₂₄", "(0,1,2)(1,0,1)₂₄"],
        ["σ² (잔차분산)", "467,463", "971,913"],
        ["log-likelihood", "−69,441", "−72,354"],
        ["AIC", "138,904", "144,723"],
        ["BIC", "138,982", "144,780"],
    ]
    add_table(s, x=0.7, y=1.7, w=6.5, h=3.6, data=table_data,
              header=True, first_col_bold=True, highlight_col=1)

    # right: insights
    add_text(s, 7.6, 1.7, 5.2, 0.4,
             "핵심 차이",
             size=18, bold=True, color=DARK)

    add_bullet(s, 7.6, 2.2, 5.4, 3.0, [
        "ΔAIC = −5,819   (경험칙 10 초과 시 결정적)",
        "σ² 비율 ÷ 2.08   → 예측 오차 표준편차 ≈ 1.44× 작음",
        "rustima : D=1 (계절차분) + 풍부한 ARMA 구조까지 탐색",
        "pmdarima : d=1 (일반차분)에 조기 수렴 → 24시간 주기를 ta·intercept가 흡수, 해석 왜곡",
    ], size=13)

    # bottom callout
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.7), Inches(5.6), Inches(12.0), Inches(1.1))
    box.fill.solid()
    box.fill.fore_color.rgb = ORANGE
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.3)
    tf.margin_top = tf.margin_bottom = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "5× 빠르고  ·  35× 가볍고  ·  AIC −5,819 만큼 더 정확하다"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT

    add_footer(s, "Same data, same exog, single Mac M-series machine")


def main():
    if not os.path.exists(CHART_UNIFIED):
        raise FileNotFoundError(f"missing chart: {CHART_UNIFIED}")

    slide_performance()
    slide_quality()

    prs.save(OUT_PATH)
    print(f"saved → {OUT_PATH}")


if __name__ == "__main__":
    main()
